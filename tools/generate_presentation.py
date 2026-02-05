from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Slide content
slides = [
    {"title": "Pedestrian Navigation with ESP32-CAM",
     "subtitle": "Real-time obstacle detection, mapping and alerts"},

    {"title": "Problem & Motivation",
     "bullets": [
         "Pedestrians face hazards (potholes, obstacles) while walking",
         "Existing navigation apps don't warn of local obstacles in real-time",
         "Low-cost camera + edge/phone detection can fill this gap"
     ]},

    {"title": "Project Objectives",
     "bullets": [
         "Real-time obstacle detection and location reporting",
         "Browser-based dashboard with live camera feed & map",
         "Low-cost hardware integration (ESP32-CAM)",
         "Trainable custom models for pothole/crack detection"
     ]},

    {"title": "System Architecture",
     "bullets": [
         "ESP32-CAM streams video (optional)",
         "Browser or Python (YOLOv8) performs detection",
         "Flask + Socket.IO server stores and broadcasts obstacles",
         "Leaflet.js map shows geo-tagged obstacles in real-time"
     ]},

    {"title": "Detection & Models",
     "bullets": [
         "Prototype: TensorFlow.js COCO-SSD in browser (fast)",
         "Server/desktop: YOLOv8 (Ultralytics) for higher accuracy",
         "Planned: Custom YOLOv8 trained for potholes/cracks",
         "Deduplication, distance-based warnings, configurable thresholds"
     ]},

    {"title": "Web Application Features",
     "bullets": [
         "Split-screen map + live camera (WebRTC or ESP32 stream)",
         "Real-time WebSocket updates and analytics dashboard",
         "Live detection panel with distance & priority warnings",
         "Permission diagnostic tool and mobile-friendly UI"
     ]},

    {"title": "ESP32-CAM Hardware",
     "bullets": [
         "AI-Thinker ESP32-CAM (OV2640) streaming at /stream",
         "FTDI programmer for uploads, 5V power, optional SD card",
         "Mounting options: dashboard, helmet, chest, handlebar",
         "Power: power bank or car USB for field deployment"
     ]},

    {"title": "Demo (How to run)",
     "bullets": [
         "Start Flask server: python web_app/server.py",
         "Open: http://localhost:5000 (or phone on same WiFi)",
         "Start camera and allow permissions (use external browser)",
         "Optional: run esp32_detector.py to process ESP32 stream"
     ]},

    {"title": "Results & Improvements",
     "bullets": [
         "Higher accuracy with 75% threshold and YOLOv8 on backend",
         "Deduplication prevents repeated alerts (30s cooldown)",
         "Distance warnings: critical <3m, warning <10m",
         "Real-time markers appear on map with user reports"
     ]},

    {"title": "Next Steps",
     "bullets": [
         "Train and deploy custom pothole/crack model (YOLOv8)",
         "Integrate ESP32-CAM as primary camera source",
         "Improve tracking to persist obstacles across frames",
         "User testing and dataset collection in target areas"
     ]},

    {"title": "Contact & Acknowledgements",
     "bullets": [
         "Your Name — Project Lead",
         "Supervisor: Prof. [Name]",
         "Repo: https://github.com/Mohammedrafiullah1928/FINAL_YEAR_PROJECT",
         "Questions? demo@youremail.com"
     ]}
]

# Styling helpers
TITLE_FONT_SIZE = Pt(36)
SUBTITLE_FONT_SIZE = Pt(18)
BULLET_FONT_SIZE = Pt(18)
ACCENT_COLOR = RGBColor(6, 90, 115)  # teal-dark
BACKGROUND_COLOR = RGBColor(245, 248, 250)  # near-white

prs = Presentation()
# Set slide dimensions if desired (default is fine)

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = slides[0]["title"]
title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR

subtitle.text = slides[0]["subtitle"]
subtitle.text_frame.paragraphs[0].font.size = SUBTITLE_FONT_SIZE
subtitle.text_frame.paragraphs[0].font.italic = True

# Content slides
for s in slides[1:]:
    slide_layout = prs.slide_layouts[1]  # title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = s["title"]
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(s.get("bullets", [])):
        p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.font.size = BULLET_FONT_SIZE
        p.font.color.rgb = RGBColor(34, 34, 34)

# Save file
out_path = "c:\\Users\\N\\Desktop\\FINAL_YEAR_PROJECT\\pedestrian-navigation-esp32cam\\FINAL_YEAR_PROJECT_presentation.pptx"
prs.save(out_path)
print(f"Presentation created: {out_path}")
